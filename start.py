import speech_recognition as sr
import pyttsx3
import webbrowser
import datetime
import os
import google.generativeai as genai
import time
import random
import win32com.client
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Pt  

def generate_slide_content(topic):
    try:
        response = model.generate_content(f"Generate a concise short note (30-50 words) about {topic}. Avoid unnecessary headings or formatting.")
        if response and hasattr(response, "text"):
            return response.text.strip()
        elif response and hasattr(response, "candidates"):
            return response.candidates[0].text.strip()
        else:
            return "Content not generated."
    except Exception as e:
        return f"Error generating content: {e}"

def say(text):
    engine = pyttsx3.init()
    engine.say(text)
    engine.runAndWait()

def takecommand():
    r = sr.Recognizer()
    
    with sr.Microphone() as source:
        print("Listening...")
        say("Listening...")
        
        r.adjust_for_ambient_noise(source)  # Helps with background noise

        while True:
            try:
                audio = r.listen(source, timeout=5)  # Listen with a timeout
                query = r.recognize_google(audio, language="en-in").lower().strip()
                
                # ✅ Handle short words
                if len(query) == 1:  # Single character words like "a", "o", etc.
                    print(f"Detected short word: {query}. Please confirm.")
                    say(f"You said {query}. Is that correct?")
                    confirmation = takecommand()
                    if "yes" in confirmation or "correct" in confirmation:
                        return query
                    else:
                        print("Let's try again...")
                        continue  # Retry input
                
                # ✅ Handle normal input
                if query:
                    print(f"User said: {query}")
                    return query

            except sr.UnknownValueError:  
                print("Sorry, I didn't catch that. Please speak clearly.")
                say("I didn't catch that. Please speak clearly.")
                continue  # Retry

            except sr.RequestError:  
                print("Speech service issue. Check your internet connection!")
                say("There is an issue with my speech recognition. Please check your internet.")
                return "error"

            except Exception as e:
                print(f"Error: {str(e)}")
                say("An error occurred. Please try again.")
                return "error"



if __name__ == '__main__':
    API_KEY = "api key"
    genai.configure(api_key=API_KEY)
    model = genai.GenerativeModel("gemini-1.5-pro-latest")
    say("hello I am aod A.I ")
    print("hell0 i am aod ai")
    time.sleep(1)
    say("what you need help for ")
    print("what you need help for ")
    print("Listening...")
    text= takecommand()
    if "game".lower() in text.lower():
        while True:
            say("lets play a game of gussing random numbers ")
            say("do you play this game say done or no")            
            print("Listening...")
            text= takecommand()
            if "no".lower() in text.lower():
                exit()
            if "done".lower() in text.lower():
                words = ["apple", "banana"]
                print(words)
                say("guess from this")
                human_score = 0
                machine_score = 0
                print("Welcome to the Word Guessing Game! ")
                say("Welcome to the Word Guessing Game! ")
                print("First to 5 points wins!\n")
                say("First to 5 points wins!\n")

                while human_score < 5 and machine_score < 5:
                    target_word = random.choice(words)
                    say("Guess the secret word from the list")
                    print(f"Words: {', '.join(words)}")
                    
                    human_guess = takecommand()
                    print(f"Your guess: {human_guess}")

                    if human_guess == target_word:
                        print("🎉 You guessed correctly! +1 point!")
                        say("You guessed correctly! +1 point!")
                        human_score += 1
                    else:
                        print(f"❌ Wrong guess! The correct word was {target_word}.")
                    machine_guess = random.choice(words)
                    print(f"🤖 Machine guesses: {machine_guess}")
                    say(f"My guess is {machine_guess}")

                    if machine_guess == target_word:
                        print("🤖 Machine guessed correctly! +1 point!")
                        say("Machine guessed correctly! +1 point!")
                        machine_score += 1
                    else:
                        print(f"🤖 Machine was wrong! The correct word was {target_word}.")
                    print(f"Score: You = {human_score} | Machine = {machine_score}\n")
                    time.sleep(1)
                if human_score == 5:
                    print("🎉🎊 Congratulations! You won the game! 🏆")
                    say("Congratulations! You won the game!")
                else:
                    print("🤖 Machine wins! Better luck next time! 😢")
                    say("Machine wins! Better luck next time!")
            if "exit".lower() in text.lower():
                say("have a nice day sir good bye")
                exit()
    
    if "work".lower() in text.lower():
        while True:
            print("Listening...")
            text= takecommand()
            if any(keyword in text.lower() for keyword in ["draft mail", "mail", "compose mail", "write email"]):
                say("Say the email address before at the rate")
                print("Tell email address before @")
                email=takecommand()
                email=email.replace(" ","")
                email=(f"{email}@gmail.com")
                say("tell the subject")
                print("tell me the subject")
                subject=takecommand()
                response = model.generate_content(f"write a message about {subject} and in should be reasonable and good without error i should be of three to four line")
                print(f"email={email} subject={subject} mess={response.text}")
                gmail_link = f"https://mail.google.com/mail/u/0/?view=cm&fs=1&to={email}&su={subject}&body={response.text}"
                webbrowser.open(gmail_link)
                
            if "word file".lower() in text.lower():
                say("Enter the topic ")
                subject = takecommand()
                response = model.generate_content(f" tkis is going to be in the word file so plz write in the proffessional and good way  '{subject}'")
                user_text = response.text
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = True  
                doc = word.Documents.Add()
                doc.Range(0, 0).Text = user_text
                
            if any(keyword in text.lower() for keyword in ["ppt", "presentation", "presentation file"]):
                say("Tell me the topic for the presentation: ")
                print("Tell me the topic")
                presentation_topic = takecommand()
                slide_topics = [
                    "Introduction",
                    "Background",
                    "Key Concepts",
                    "Importance",
                    "Challenges",
                    "Solutions",
                    "Case Studies",
                    "Future Trends",
                    "Conclusion"
                ]
                ppt = Presentation()
                title_slide_layout = ppt.slide_layouts[0]  
                title_slide = ppt.slides.add_slide(title_slide_layout)
                title = title_slide.shapes.title
                title.text = presentation_topic
                for slide_topic in slide_topics:
                    slide_layout = ppt.slide_layouts[1]  
                    slide = ppt.slides.add_slide(slide_layout)
                    title = slide.shapes.title
                    title.text = slide_topic
                    slide_content = generate_slide_content(f"{presentation_topic} - {slide_topic}")
                    content_placeholder = slide.placeholders[1]
                    content_placeholder.text = slide_content  
                    for paragraph in content_placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(22)
                ppt_filename = "generated_presentation.pptx"
                ppt.save(ppt_filename)
                os.system(f'start {ppt_filename}')
                print(f"PowerPoint presentation '{ppt_filename}' generated successfully!")

            
            
            if "exit".lower() in text.lower():
                say("have a nice day sir good bye")
                exit()
    print("not understand")
